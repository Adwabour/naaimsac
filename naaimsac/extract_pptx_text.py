import sys
from pptx import Presentation

if len(sys.argv) != 2:
    print("Usage: python3 extract_pptx_text.py <pptx_path>")
    sys.exit(1)

pptx_path = sys.argv[1]
pres = Presentation(pptx_path)

for i, slide in enumerate(pres.slides):
    print(f"\n--- Slide {i+1} ---")
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            print(shape.text)